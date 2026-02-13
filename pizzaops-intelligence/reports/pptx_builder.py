"""
PowerPoint report generation for PizzaOps Intelligence.
Creates professional slide decks using python-pptx.
"""

import pandas as pd
from typing import Dict, List, Optional, Any
from datetime import datetime
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True

    # Color definitions - moved inside the try block
    COLORS = {
        'primary': RGBColor(255, 107, 53),      # Orange
        'primary_dark': RGBColor(229, 90, 43),
        'success': RGBColor(16, 185, 129),
        'warning': RGBColor(245, 158, 11),
        'danger': RGBColor(239, 68, 68),
        'dark': RGBColor(14, 17, 23),
        'gray': RGBColor(148, 163, 184),
        'light': RGBColor(250, 250, 250),
        'white': RGBColor(255, 255, 255),
    }
except ImportError:
    PPTX_AVAILABLE = False
    # Define a default COLORS dictionary if pptx is not available
    # to avoid NameError if other parts of the module try to access it
    COLORS = {
        'primary': (255, 107, 53),
        'primary_dark': (229, 90, 43),
        'success': (16, 185, 129),
        'warning': (245, 158, 11),
        'danger': (239, 68, 68),
        'dark': (14, 17, 23),
        'gray': (148, 163, 184),
        'light': (250, 250, 250),
        'white': (255, 255, 255),
    }


def create_presentation() -> Any:
    """Create a new presentation."""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(
    prs: Any,
    title: str,
    subtitle: str = "",
    date_range: str = ""
) -> None:
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11.333), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4),
            Inches(11.333), Inches(0.75)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Date range
    if date_range:
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.75),
            Inches(11.333), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Report Period: {date_range}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Branding
    brand_box = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(11.333), Inches(0.5)
    )
    tf = brand_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PizzaOps Intelligence by JLWanalytics"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER


def add_kpi_slide(
    prs: Any,
    title: str,
    kpis: List[Dict]
) -> None:
    """Add a slide with KPI cards."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.75)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # KPI cards
    card_width = Inches(2.8)
    card_height = Inches(1.5)
    start_x = Inches(0.75)
    start_y = Inches(1.5)
    gap = Inches(0.3)

    for i, kpi in enumerate(kpis[:4]):
        x = start_x + (i * (card_width + gap))

        # Determine color
        status = kpi.get('status', 'neutral')
        if status == 'good':
            color = COLORS['success']
        elif status == 'warning':
            color = COLORS['warning']
        elif status == 'danger':
            color = COLORS['danger']
        else:
            color = COLORS['gray']

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, start_y,
            card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(240, 240, 240)
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Label
        label_box = slide.shapes.add_textbox(
            x + Inches(0.15), start_y + Inches(0.15),
            card_width - Inches(0.3), Inches(0.4)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get('label', '')
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']

        # Value
        value_box = slide.shapes.add_textbox(
            x + Inches(0.15), start_y + Inches(0.55),
            card_width - Inches(0.3), Inches(0.75)
        )
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(kpi.get('value', ''))
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color


def add_bullet_slide(
    prs: Any,
    title: str,
    bullets: List[str],
    subtitle: str = ""
) -> None:
    """Add a slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.75)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1),
            Inches(12.333), Inches(0.5)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray']

    # Bullets
    start_y = Inches(1.75) if subtitle else Inches(1.25)
    bullet_box = slide.shapes.add_textbox(
        Inches(0.75), start_y,
        Inches(11.833), Inches(5)
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:8]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(12)


def add_table_slide(
    prs: Any,
    title: str,
    df: pd.DataFrame
) -> None:
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.75)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Table
    rows = min(len(df) + 1, 10)  # +1 for header, max 10 rows
    cols = min(len(df.columns), 6)  # Max 6 columns

    table = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.25),
        Inches(12.333), Inches(5)
    ).table

    # Set column widths
    col_width = Inches(12.333 / cols)
    for i in range(cols):
        table.columns[i].width = int(col_width)

    # Header row
    for j, col_name in enumerate(df.columns[:cols]):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']

        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']

    # Data rows
    for i, (_, row) in enumerate(df.head(rows - 1).iterrows()):
        for j, col in enumerate(df.columns[:cols]):
            cell = table.cell(i + 1, j)
            value = row[col]
            if isinstance(value, float):
                value = f"{value:.2f}"
            cell.text = str(value)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark']


def add_insight_slide(
    prs: Any,
    title: str,
    main_stat: str,
    main_label: str,
    description: str
) -> None:
    """Add a slide highlighting a key insight."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.75)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Main stat
    stat_box = slide.shapes.add_textbox(
        Inches(2), Inches(2),
        Inches(9.333), Inches(1.5)
    )
    tf = stat_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_stat
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(2), Inches(3.5),
        Inches(9.333), Inches(0.5)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_label
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5),
        Inches(11.333), Inches(2)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER


def generate_executive_presentation(
    kpis: Dict,
    area_metrics: pd.DataFrame,
    bottlenecks: List[Dict],
    complaint_analysis: Dict,
    recommendations: List[str],
    date_range: str = ""
) -> bytes:
    """
    Generate executive presentation.

    Args:
        kpis: Overview KPIs
        area_metrics: Area performance
        bottlenecks: Bottleneck list
        complaint_analysis: Complaint analysis results
        recommendations: Recommendation strings
        date_range: Report period

    Returns:
        PPTX bytes
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    prs = create_presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        title="PizzaOps Performance Report",
        subtitle="Executive Summary",
        date_range=date_range
    )

    # Slide 2: KPIs
    kpi_list = [
        {
            'label': 'Total Orders',
            'value': f"{kpis.get('total_orders', 0):,}",
            'status': 'neutral'
        },
        {
            'label': 'On-Time %',
            'value': f"{kpis.get('on_time_pct', 0):.1f}%",
            'status': 'good' if kpis.get('on_time_pct', 0) >= 85 else 'warning'
        },
        {
            'label': 'Complaint Rate',
            'value': f"{kpis.get('complaint_rate', 0):.1f}%",
            'status': 'good' if kpis.get('complaint_rate', 0) <= 5 else 'danger'
        },
        {
            'label': 'Avg Delivery',
            'value': f"{kpis.get('avg_delivery_time', 0):.1f} min",
            'status': 'good' if kpis.get('avg_delivery_time', 0) <= 25 else 'warning'
        }
    ]
    add_kpi_slide(prs, "Key Performance Indicators", kpi_list)

    # Slide 3: Key Insight (Complaint Root Cause)
    if complaint_analysis:
        on_time_pct = complaint_analysis.get('on_time_complaint_pct', 0)
        add_insight_slide(
            prs,
            title="Key Insight: Complaint Root Cause",
            main_stat=f"{on_time_pct:.0f}%",
            main_label="of complaints from ON-TIME deliveries",
            description="This means complaints are NOT just about speed. "
                       "Many customers complain about quality issues (cold food, wrong orders) "
                       "even when delivery is fast. Focus on oven temperature and order accuracy."
        )

    # Slide 4: Area Performance
    if len(area_metrics) > 0:
        add_table_slide(prs, "Delivery Performance by Area", area_metrics)

    # Slide 5: Recommendations
    add_bullet_slide(
        prs,
        title="Recommendations",
        bullets=recommendations[:6],
        subtitle="Prioritized actions based on data analysis"
    )

    # Slide 6: Next Steps
    next_steps = [
        "Review oven preheating protocols to reduce cold food complaints",
        "Implement order verification checklist to reduce wrong order issues",
        "Optimize delivery routes for Areas E and C",
        "Adjust staffing during peak hours (11-14, 17-21)",
        "Track KPIs weekly to monitor improvement"
    ]
    add_bullet_slide(prs, "Next Steps", next_steps)

    # Save to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
